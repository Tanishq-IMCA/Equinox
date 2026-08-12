from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from docx import Document
from docx.shared import Inches as DInches, Pt as DPt, RGBColor as DRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "deliverables"
OUT.mkdir(exist_ok=True)

NAVY = RGBColor(8, 12, 24)
PANEL = RGBColor(17, 24, 39)
PANEL_2 = RGBColor(24, 31, 50)
WHITE = RGBColor(245, 247, 255)
MUTED = RGBColor(168, 176, 196)
FAINT = RGBColor(102, 113, 138)
PURPLE = RGBColor(183, 117, 255)
GREEN = RGBColor(66, 245, 155)
CYAN = RGBColor(105, 231, 255)
YELLOW = RGBColor(250, 204, 21)
RED = RGBColor(255, 95, 103)

W, H = 13.333, 7.5


def set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, size=16, color=WHITE, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.02, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(2)
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def line(slide, x1, y1, x2, y2, color=FAINT, width=1.2):
    shape = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    return shape


def footer(slide, index, label="EQUINOX  ·  ENERGY-AWARE WORKLOAD ORCHESTRATION"):
    line(slide, 0.55, 7.05, 12.78, 7.05, RGBColor(40, 48, 70), 0.7)
    add_text(slide, label, 0.58, 7.12, 9.5, 0.18, 7.5, FAINT, font="Aptos", margin=0)
    add_text(slide, f"{index:02d} / 21", 11.75, 7.12, 1.0, 0.18, 7.5, FAINT, align=PP_ALIGN.RIGHT, margin=0)


def chrome(slide, index, section="TECHNICAL PRESENTATION"):
    add_text(slide, section.upper(), 0.58, 0.35, 6, 0.2, 8, PURPLE, bold=True, margin=0)
    add_text(slide, "EQUINOX", 11.5, 0.35, 1.25, 0.2, 8, WHITE, bold=True, align=PP_ALIGN.RIGHT, margin=0)
    footer(slide, index)


def heading(slide, kicker, title, subtitle=None, index=None):
    add_text(slide, kicker.upper(), 0.65, 0.85, 5.8, 0.22, 9, PURPLE, bold=True, margin=0)
    add_text(slide, title, 0.62, 1.12, 11.9, 0.72, 31, WHITE, bold=True, font="Aptos Display", margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.66, 1.92, 11.2, 0.46, 13, MUTED, margin=0)
    if index is not None:
        chrome(slide, index)


def card(slide, x, y, w, h, label, title, body, accent=PURPLE, body_size=11.5):
    rect(slide, x, y, w, h, PANEL, RGBColor(42, 51, 76), True)
    rect(slide, x, y, 0.06, h, accent, accent)
    add_text(slide, label.upper(), x + 0.2, y + 0.18, w - 0.35, 0.2, 8, accent, bold=True, margin=0)
    add_text(slide, title, x + 0.2, y + 0.48, w - 0.35, 0.46, 16, WHITE, bold=True, margin=0)
    add_text(slide, body, x + 0.2, y + 1.02, w - 0.35, h - 1.15, body_size, MUTED, margin=0)


def bullets(slide, items, x, y, w, h, size=13, color=MUTED, accent=GREEN, gap=0.28):
    yy = y
    for item in items:
        rect(slide, x, yy + 0.1, 0.07, 0.07, accent, accent)
        add_text(slide, item, x + 0.18, yy, w - 0.2, gap, size, color, margin=0)
        yy += gap


def section_slide(prs, index, kicker, title, subtitle, cards):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    heading(slide, kicker, title, subtitle, index)
    n = len(cards)
    gap = 0.22
    margin = 0.65
    width = (12.0 - gap * (n - 1)) / n
    for i, data in enumerate(cards):
        card(slide, margin + i * (width + gap), 2.7, width, 3.6, *data)
    return slide


def build_deck():
    source = ROOT / "attached_assets" / "AURA_Technical_Presentation_1786515526385.pptx"
    prs = Presentation(str(source))
    while len(prs.slides):
        rid = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[-1]
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    rect(s, 0.65, 0.6, 0.12, 5.8, PURPLE, PURPLE)
    add_text(s, "TECHNICAL PRESENTATION  ·  2026", 1.02, 0.72, 6, 0.25, 10, PURPLE, bold=True, margin=0)
    add_text(s, "EQUINOX", 1.0, 1.55, 8.8, 0.9, 48, WHITE, True, "Aptos Display", margin=0)
    add_text(s, "Autonomous energy-aware\nworkload orchestration", 1.02, 2.58, 8.2, 1.05, 25, CYAN, True, "Aptos Display", margin=0)
    add_text(s, "A live control system for moving work, balancing heat, and turning idle infrastructure into capacity instead of waste.", 1.05, 4.05, 6.7, 0.8, 15, MUTED, margin=0)
    rect(s, 9.05, 1.45, 3.25, 3.9, PANEL, RGBColor(53, 42, 79), True)
    add_text(s, "THE DEMO", 9.35, 1.78, 2.5, 0.2, 9, GREEN, True, margin=0)
    add_text(s, "60", 9.32, 2.2, 1.3, 0.65, 38, WHITE, True, "Aptos Display", margin=0)
    add_text(s, "simulated racks\nin the current UI", 10.42, 2.28, 1.45, 0.5, 11, MUTED, margin=0)
    add_text(s, "8", 9.32, 3.15, 1.3, 0.65, 38, WHITE, True, "Aptos Display", margin=0)
    add_text(s, "workload\ncategories", 10.42, 3.23, 1.45, 0.5, 11, MUTED, margin=0)
    add_text(s, "10 s", 9.32, 4.1, 1.5, 0.65, 32, PURPLE, True, "Aptos Display", margin=0)
    add_text(s, "slow intake\ncadence", 10.42, 4.2, 1.45, 0.5, 11, MUTED, margin=0)
    add_text(s, "PRESENTED BY", 1.05, 6.25, 1.5, 0.2, 8, FAINT, True, margin=0)
    add_text(s, "Tanishq Giri  ·  Aaroh Dharmadhikari", 2.35, 6.25, 5.8, 0.2, 9, WHITE, True, margin=0)
    footer(s, 1)

    # 2 contents
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); heading(s, "CONTENTS", "The room, in motion.", "A fast path from why this matters to how the prototype actually works.", 2)
    items = ["01  General start", "02  Problem + relevance", "03  Innovation + creativity", "04  Architecture + APIs", "05  Complexity + execution", "06  Working prototype", "07  UI/UX system", "08  Impact + scalability", "09  Collaboration + openness", "10  The wow factor"]
    for i, item in enumerate(items):
        x = 0.85 + (i // 5) * 6.0; y = 2.65 + (i % 5) * 0.72
        add_text(s, item, x, y, 5.25, 0.3, 15, WHITE if i % 2 == 0 else MUTED, bold=i % 2 == 0, margin=0)
        line(s, x, y + 0.42, x + 5.15, y + 0.42, RGBColor(42, 51, 76), 0.6)

    section_slide(prs, 3, "01 · GENERAL START", "Compute is physical.", "Equinox treats a cluster as a living system: tasks arrive, resources heat up, work moves, and capacity changes state.", [
        ("THE IDEA", "Make compute work smarter.", "A workload is not just a request. It has a thermal cost, a power cost, a memory footprint, and a lifetime. Equinox makes those costs visible and actionable.", PURPLE),
        ("THE TARGET", "One better cluster state.", "The engine continuously searches for safer, denser placements instead of accepting the first available node forever.", GREEN),
        ("THE DEMO", "A room you can operate.", "The dashboard turns allocation into a visual control surface: racks, lights, telemetry panels, task lists, logs, and power actions.", CYAN),
    ])
    section_slide(prs, 4, "02 · PROBLEM STATEMENT", "The problem is not “more servers.”", "It is unmanaged state: uneven workload density, hidden thermal pressure, and idle capacity left powered on.", [
        ("OPERATIONS", "Static placement wastes headroom.", "A task can be technically running while the room is operationally inefficient. One node saturates while peers sit cold and empty.", RED),
        ("ENERGY", "Idle capacity still costs.", "Powering every rack for a burst that has already ended makes the infrastructure pay for yesterday’s demand.", YELLOW),
        ("VISIBILITY", "Telemetry without action is a report.", "CPU and temperature numbers matter only when they influence placement, migration, shutdown, or restart decisions.", PURPLE),
    ])
    section_slide(prs, 5, "02 · RELEVANCE", "Why this matters now.", "AI workloads, inference, games, storage, and services do not share the same resource shape — or the same heat curve.", [
        ("REAL SYSTEMS", "Mixed workloads are the default.", "The task catalog represents operating-system work, background services, ML, LLM execution, gaming, hosting, disk I/O, and general tasks.", GREEN),
        ("RESOURCE PRESSURE", "Every resource has a ceiling.", "The prototype models CPU, RAM, GPU, VRAM, power, and temperature against per-node envelopes instead of a single generic “load” score.", CYAN),
        ("HUMAN VALUE", "Operators need a readable state.", "Equinox makes the invisible tradeoff legible: accept, move, consolidate, power down, or keep capacity online.", PURPLE),
    ])
    section_slide(prs, 6, "03 · INNOVATION", "The innovation is choreography.", "Equinox turns capacity management into a visible sequence of decisions instead of a silent threshold alarm.", [
        ("POWER-AWARE", "Lights out is a feature.", "In Autonomous mode, work is drained, nodes transition gradually, and only demand reactivates capacity. The room itself communicates the algorithm.", PURPLE),
        ("THERMAL-AWARE", "Heat has memory.", "Temperature eases toward a target derived from task heat and node state. The visualization avoids unrealistic instant jumps.", RED),
        ("CONTROL LOOP", "Allocation never really stops.", "A slow task intake cadence, finite task lifetimes, migration, completion, consolidation, and reactivation create a continuous simulation.", GREEN),
    ])

    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); heading(s, "04 · SYSTEM ARCHITECTURE", "A small stack with a clear control loop.", "Next.js owns the surface and internal API routes; the browser simulation owns the live runtime state.", 7)
    layers = [("BROWSER SURFACE", "Landing · About · Login · Dashboard", PURPLE), ("INTERACTION", "React state · Three.js scene · telemetry panels", CYAN), ("SIMULATION CORE", "useSimulation() · task intake · placement · power state", GREEN), ("INTERNAL API", "GET /api/tasks · GET /api/powerwall", YELLOW), ("ASSETS", "CSV task library · GLB rack · GLB Powerwall · fonts", RED)]
    for i, (name, body, color) in enumerate(layers):
        y = 2.35 + i * 0.75
        rect(s, 1.0 + i * 0.24, y, 10.7 - i * 0.48, 0.52, PANEL_2, color, True)
        add_text(s, name, 1.25 + i * 0.24, y + 0.1, 2.2, 0.2, 9, color, True, margin=0)
        add_text(s, body, 3.4 + i * 0.24, y + 0.1, 7.6, 0.2, 12, WHITE, margin=0)
        if i < len(layers) - 1: line(s, 6.4, y + 0.52, 6.4, y + 0.73, color, 1.3)
    add_text(s, "The important design choice: the demo does not pretend to be a distributed production scheduler. It is a transparent, deterministic prototype of the decisions a scheduler would make.", 1.0, 6.25, 10.8, 0.42, 12, MUTED, italic=True, margin=0)
    footer(s, 7)

    section_slide(prs, 8, "04 · DATA ENGINE", "The task library becomes live work.", "The CSV is not decoration: it is parsed into typed templates and becomes the source of resource-aware workloads.", [
        ("SOURCE", "agenticwork/tasks_preview.csv", "The library contains 400 task records across eight categories with unique IDs, believable names, intensity tiers, and range-based resource fields.", CYAN),
        ("NORMALIZE", "Average the ranges.", "The route handler extracts numeric ranges and averages their endpoints. The browser receives stable numeric values for simulation.", GREEN),
        ("EXECUTE", "Give every task a lifetime.", "makeTask() creates a unique runtime UID and a bounded duration; heavy categories receive longer minimum lifetimes.", PURPLE),
    ])
    section_slide(prs, 9, "04 · API + SPONSOR TECH", "Small APIs. Honest boundaries.", "The current build uses internal Next.js route handlers and open-source runtime libraries; no external sponsor integration is connected.", [
        ("TASK API", "GET /api/tasks", "Reads the local CSV, parses rows, converts ranges to numeric averages, and returns JSON consumed by useSimulation().", GREEN),
        ("ASSET API", "GET /api/powerwall", "Streams the attached Powerwall GLB with the model/gltf-binary content type so Three.js can mount it in the room.", CYAN),
        ("SPONSOR TECH", "Not currently connected.", "Three.js, React, Next.js, and Framer Motion are open-source dependencies. A future sponsor API can feed real telemetry behind the same typed boundary.", PURPLE),
    ])
    section_slide(prs, 10, "05 · TECHNICAL COMPLEXITY", "The hard part is state, not screens.", "The prototype coordinates independent lifecycles: task completion, migration, node power actions, thermal easing, and UI rendering.", [
        ("PLACEMENT", "Multi-resource admission.", "canAccept() checks CPU, RAM, GPU, VRAM, power, and a temperature constraint against an 80% operating envelope.", GREEN),
        ("MIGRATION", "Move without overload.", "Transfers are blocked when a peer cannot accept the full task set. Manual termination marks a task as transferring before it disappears.", PURPLE),
        ("POWER", "Staggered state transitions.", "Starting and stopping are asynchronous, sequentially delayed, and reflected in both telemetry state and 3D rack lights.", YELLOW),
    ])
    section_slide(prs, 11, "05 · EXECUTION", "What happens every second.", "The one-second loop updates temperatures, decrements task lifetimes, removes completed work, and decides whether intake or consolidation is due.", [
        ("TICK", "Telemetry evolves.", "Each node moves toward a target temperature. Offline nodes cool; active nodes respond to their current task heat.", RED),
        ("INTAKE", "Every 10 seconds.", "The prototype intentionally uses a slow arrival cadence. In normal mode it assigns directly; in Autonomous mode it queues and dispatches.", CYAN),
        ("CONTROL", "Capacity follows demand.", "Autonomous mode activates a node only when queued work cannot fit, then consolidates empty or movable nodes back down.", GREEN),
    ])
    section_slide(prs, 12, "06 · FUNCTIONALITY", "A working prototype, not a slideware promise.", "The dashboard has a login gate, a live 3D room, node hover telemetry, cluster task management, power controls, autonomous mode, and logs.", [
        ("OPERATE", "Overview.", "Rotate and zoom the rack room, hover a rack, inspect live metrics, read cluster summaries, toggle Autonomous mode, and power a cluster sequentially.", PURPLE),
        ("INSPECT", "Live Workflow.", "Select a cluster, review its process list, inspect all node telemetry, and request a task transfer/termination.", CYAN),
        ("AUDIT", "Logs.", "Engine, autonomous, power, transfer, completion, and notice events form a readable narrative of what the control loop did.", GREEN),
    ])
    section_slide(prs, 13, "06 · AUTONOMOUS CONTROL", "The loop is designed to stay alive.", "Autonomous mode first drains and powers down capacity, then becomes a demand-driven allocator with continuous slow intake.", [
        ("ENGAGE", "Preserve then drain.", "Current workloads move into a pending queue. Nodes stop sequentially so the room visibly reaches a lower-power baseline.", YELLOW),
        ("ALLOCATE", "Queue → fit → activate.", "Queued work is placed onto an accepting node. If none can fit, the next offline node starts and dispatch resumes.", GREEN),
        ("RECOVER", "Consolidate after work ends.", "The runtime moves compatible tasks together and schedules empty nodes for shutdown, preventing green but idle racks from accumulating.", PURPLE),
    ])
    section_slide(prs, 14, "07 · UI / UX", "A control room with a calm signal.", "The UI uses dark glass panels, purple/cyan accents, readable operational states, and a deliberate hierarchy from room-level to node-level detail.", [
        ("ORIENTATION", "Three levels of zoom.", "Landing page explains the idea. Overview shows the room. Live Workflow exposes cluster and node operations. Logs explain the timeline.", PURPLE),
        ("FEEDBACK", "State is visual and textual.", "Rack lights, colored labels, capacity bars, task counts, notices, and log events reinforce the same state through different channels.", GREEN),
        ("ACCESS", "Readable under pressure.", "The dashboard uses explicit labels, disabled states, status colors, and enlarged small text so operational metadata is not treated as decoration.", CYAN),
    ])
    section_slide(prs, 15, "08 · IMPACT", "Less waste. More usable compute.", "The immediate value is not a magical energy claim; it is a measurable control strategy that makes energy-aware placement understandable.", [
        ("EFFICIENCY", "Consolidate idle work.", "When compatible tasks share a node, empty servers can move toward offline, reducing modeled power draw and visual clutter.", GREEN),
        ("RESILIENCE", "Keep headroom visible.", "The 80% admission target leaves room for variance, while high-state visuals communicate when the room is approaching pressure.", YELLOW),
        ("OPERABILITY", "Make decisions explainable.", "Each transfer, activation, shutdown, and completion is logged so an operator can answer “why did this node change?”", PURPLE),
    ])
    section_slide(prs, 16, "08 · SCALABILITY", "Scale the boundary, not the confusion.", "The prototype can grow from a browser simulation into a real scheduler by replacing data sources behind the same control concepts.", [
        ("NODES", "30 → 60 → real clusters.", "The original brief targets 30 nodes; the current visual demo renders 60 nodes across six clusters. The next step is config-driven topology.", CYAN),
        ("DATA", "CSV → stream.", "Replace the local task route with a queue, telemetry broker, or scheduler adapter while retaining typed task templates and admission checks.", PURPLE),
        ("CONTROL", "Rules → model assist.", "A learned scoring model can rank placements, but hard resource and safety constraints should remain deterministic guardrails.", GREEN),
    ])
    section_slide(prs, 17, "09 · TEAM COLLABORATION", "Two people. One system.", "The split follows the product’s two halves: the infrastructure experience and the intelligence/data layer.", [
        ("TANISHQ", "Backend + frontend lead.", "Owns the Next.js application shell, dashboard UX, Three.js room, navigation, telemetry surfaces, simulation integration, visual system, and end-to-end product assembly.", PURPLE),
        ("AAROH", "Library + API + intelligence.", "Owns task-library creation and curation, API development for task/model delivery, data-science framing, workload profiles, and small validation/documentation improvements.", GREEN),
        ("TOGETHER", "Integrate at typed seams.", "Agree on the task schema, resource units, node bounds, event language, and demo narrative so data and UI evolve without breaking the control story.", CYAN),
    ])
    section_slide(prs, 18, "09 · OPEN SOURCE", "Open by contribution, honest by default.", "The repository currently has no explicit LICENSE file, so open-source availability is a decision to make before public release.", [
        ("TODAY", "Dependencies are open source.", "Next.js, React, Three.js, Framer Motion, TypeScript, and their transitive ecosystem are used as open-source building blocks.", GREEN),
        ("RECOMMEND", "Add a project license.", "MIT is a practical default for the application code; keep asset licenses and any generated task data provenance documented separately.", PURPLE),
        ("CONTRIBUTE", "Make the seams welcoming.", "Publish the schema, add setup instructions, provide fixture data, label simulation vs production behavior, and accept improvements to schedulers, adapters, and UI.", CYAN),
    ])
    section_slide(prs, 19, "10 · WOW FACTOR", "The room does not just glow. It decides.", "Equinox’s epic moment is the Autonomous transition: a busy room drains, lights go out, demand returns, and capacity wakes back up only when it has a reason.", [
        ("CINEMATIC", "Power lines tell the story.", "The Three.js scene turns abstract allocation into visible energy paths, rack states, pulsing supply, and a Powerwall asset.", CYAN),
        ("ALIVE", "Tasks have a heartbeat.", "Work arrives slowly, runs for a bounded lifetime, completes, transfers, and changes the thermal state around it.", GREEN),
        ("MEMORABLE", "The operator sees cause and effect.", "Toggle Autonomous and the room performs the algorithm in front of you. It is both a demo and a mental model for energy-aware scheduling.", PURPLE),
    ])
    section_slide(prs, 20, "11 · DEFENSE CARD", "The 60-second explanation.", "If the room is moving fast, use this answer: Equinox is a browser-based prototype of energy-aware workload orchestration.", [
        ("PROBLEM", "Clusters waste energy when workloads are placed without thermal or power context.", "Equinox models mixed tasks and node bounds so placement becomes a multi-resource decision.", RED),
        ("SOLUTION", "A continuous control loop.", "Tasks arrive every 10 seconds, live tasks consume resources, compatible workloads consolidate, and idle nodes power down.", GREEN),
        ("PROOF", "Working product surface.", "A real Next.js dashboard renders a 3D room, live telemetry, task operations, API-fed workload data, autonomous transitions, and logs.", PURPLE),
    ])
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_text(s, "THE NEXT STATE.", 0.75, 1.15, 8, 0.4, 13, PURPLE, True, margin=0)
    add_text(s, "Compute that\nknows when to move.", 0.72, 1.75, 9.5, 1.5, 42, WHITE, True, "Aptos Display", margin=0)
    add_text(s, "EQUINOX", 0.78, 4.25, 3.8, 0.5, 24, CYAN, True, "Aptos Display", margin=0)
    add_text(s, "Tanishq Giri  ·  Aaroh Dharmadhikari\nAutonomous energy-aware workload orchestration", 0.82, 5.1, 6.7, 0.7, 14, MUTED, margin=0)
    rect(s, 9.15, 1.25, 2.8, 3.95, PANEL, RGBColor(53, 42, 79), True)
    for i, c in enumerate([PURPLE, CYAN, GREEN, YELLOW, RED]):
        rect(s, 9.55, 1.75 + i * 0.53, 1.95, 0.07, c, c, True)
    add_text(s, "MAKE COMPUTE\nWORK SMARTER.", 9.55, 4.35, 2.0, 0.5, 12, WHITE, True, margin=0)
    footer(s, 21)
    out = OUT / "Equinox_Technical_Presentation.pptx"
    prs.save(str(out))
    return out


SECTIONS = [
    ("Executive quick read", "Equinox is a browser-based prototype of an autonomous, energy-aware workload orchestrator. It simulates mixed workloads entering a clustered environment, calculates multi-resource pressure, migrates or consolidates work when safe, and powers down capacity that no longer needs to be online. The product is intentionally visual: the 3D room, rack lights, telemetry panels, process lists, autonomous transitions, and logs all explain the control loop in motion.",
     ["The current code is a Next.js 14 application with React and a client-side simulation runtime.", "The current dashboard renders 60 simulated racks in six clusters; the original history brief describes 30 nodes. This is documented as a prototype discrepancy, not hidden.", "Workload intake is intentionally slow: one new autonomous task is generated approximately every 10 seconds after Autonomous mode is ready.", "The current control engine is deterministic and rule-based. It is ML-ready in its separation of task data, placement scoring, and hard admission constraints, but it does not call a trained ML model today.", "There are two internal API routes: /api/tasks for the CSV-backed task catalog and /api/powerwall for streaming the GLB visual asset.", "The repository currently has no explicit LICENSE file. Open-source release therefore requires an explicit license decision and asset-provenance review." ]),
    ("1. Project overview and mission", "Equinox exists to make compute placement respond to physical reality. A conventional demo can show CPU bars; Equinox tries to show the system-level consequence of those bars: a task creates heat, consumes energy, occupies memory, changes placement options, and eventually changes which machines need to be powered. The system’s mission is to make those decisions visible and controllable.",
     ["The product uses an energy-aware vocabulary: workload, node, cluster, capacity, thermal target, admission envelope, migration, consolidation, activation, shutdown, and event log.", "The interface is designed as an operational control room rather than a marketing-only page. The landing page sets the frame, the About page explains the system and team, and the dashboard provides live interaction.", "The prototype is deliberately honest about simulation boundaries. It does not claim to control real servers, issue real power commands, or forecast production energy savings.", "The history brief established per-node reference limits of 120 GB RAM, 60 GB VRAM, 100°C thermal cap, and 19.2 kW power envelope. The simulation normalizes values against these bounds." ]),
    ("2. Problem statement and relevance", "Data centers and clustered compute systems do not fail only because they run out of CPU. They become inefficient when work is scattered, high-power tasks cluster badly, thermal headroom disappears, or capacity remains powered even though the workload has moved elsewhere. Mixed workloads make this harder because an LLM batch, disk repair, browser worker, and game server have radically different resource profiles.",
     ["Static placement is easy to implement but poor at adapting to changing demand. It can leave one node hot while other nodes are underused.", "A single utilization percentage hides bottlenecks. CPU can be low while VRAM, power, or thermal conditions are the real constraint.", "Idle servers create a second inefficiency: they are available, but availability is not free. Autonomous mode makes capacity itself a managed state.", "The relevance extends to AI inference, model training, GPU pools, game hosting, storage services, CI runners, and mixed enterprise environments.", "Equinox’s contribution is a working explanatory model: it connects a task’s resource shape to node state, then exposes the decision sequence to an operator." ]),
    ("3. Innovation and creativity", "The creative idea is to choreograph capacity management. Instead of treating power control as a hidden backend action, Equinox turns it into an observable performance: workloads drain, rack lights change, power paths remain visible, and demand wakes capacity only when placement requires it.",
     ["The Autonomous transition creates a strong before/after state: a busy room is consolidated and powered down, then new work reactivates only the capacity required to accept it.", "Thermal values ease toward workload-derived targets rather than jumping instantly. This makes the demo feel like an infrastructure system with inertia.", "The 3D room is not a decorative scene. Rack states are synchronized with telemetry, offline racks lose their lights, labels change state color, and power pulses are shown along the supply path.", "The task library uses believable operational names across eight categories, which makes the simulation easier to explain than abstract Task 01 / Task 02 records.", "The architecture leaves a clean extension point for an ML ranker: the model can recommend an ordering, while deterministic admission checks continue to protect safety." ]),
    ("4. Architecture overview", "The application is a single Next.js project with a browser runtime. The frontend owns the interactive simulation because the product is a demonstrator; route handlers provide typed, local data boundaries. The system is small enough to understand end to end and structured enough to evolve.",
     ["App routes: app/page.tsx is the landing page, app/about/page.tsx is the product/team page, app/login/page.tsx is the local demo gate, and app/dashboard/page.tsx is the control surface.", "Simulation core: app/dashboard/simulation.ts defines task templates, runtime nodes, telemetry math, placement checks, task lifecycle, autonomous dispatch, consolidation, and power transitions.", "API routes: app/api/tasks/route.ts parses the task CSV and returns JSON; app/api/powerwall/route.ts streams the attached Powerwall GLB model.", "Assets: agenticwork/tasks_preview.csv supplies the catalog; public/data_center_rack.glb supplies the rack model; the Powerwall remains under attached_assets and is served through its route; Gotham fonts and profile images are in public.", "Runtime stack: Next.js 14.2, React 18, TypeScript, Three.js, and Framer Motion. No database, queue, cloud telemetry service, or external sponsor connector is currently connected." ]),
    ("5. Codebase map", "The codebase is intentionally compact. The key engineering story is concentrated in a few files rather than distributed across a large backend.",
     ["app/page.tsx: client-rendered landing page with the Equinox hero, glitch reveal component, ticker, product feature cards, and navigation.", "app/about/page.tsx: product brief, reference bounds, and the team cards for Tanishq and Aaroh.", "app/dashboard/page.tsx: client dashboard, route guard, Three.js scene lifecycle, rack and Powerwall loading, navigation, panels, cluster views, and event views.", "app/dashboard/simulation.ts: the state machine and control loop. It is the most important technical file because it converts templates into live tasks and nodes.", "app/api/tasks/route.ts: CSV parser, range averaging, typed response, and cache header.", "app/api/powerwall/route.ts: binary asset response with model/gltf-binary content type.", "app/globals.css and app/overrides.css: shared visual system, page layout, dashboard panel styling, typography, fonts, scrollbar treatment, and responsive rules.", "package.json and tsconfig.json: runtime dependencies, dev command, and TypeScript boundaries. The reffiles directory is excluded from the type scan because it contains imported reference material." ]),
    ("6. Data engine and task library", "The task library is the source of variety. It contains 400 entries across eight categories: Operating System, Background services, Machine learning, LLM model execution, Gaming, Hosting, Disk read/write, and General tasks. Each row carries a unique Task_ID, believable name, intensity tier, and range-based resource fields.",
     ["The CSV is intentionally range-based. A row can say 5–15 rather than pretending that every run of a task has identical resource usage.", "app/api/tasks/route.ts parses CSV lines with a small quote-aware parser. It maps the source column names to the TaskTemplate shape expected by the simulation.", "averageRange() extracts numeric values. A single number is returned directly; a two-number range becomes the mean. This turns a human-readable catalog into a stable simulation input.", "The response is cached for one hour because the task catalog is static during a demo session. A production version would likely version the library or stream tasks from a scheduler.", "A fallbackTasks list inside simulation.ts keeps the UI alive if the CSV API is unavailable. This is a safe local baseline, not a silent production data substitute." ]),
    ("7. Runtime simulation and control logic", "useSimulation() is a React hook backed by refs for mutable runtime state and a redraw counter for UI refreshes. The runtime contains nodes with tasks, offline state, powerAction, and temperature. The hook exposes telemetry, logs, notices, task termination, cluster power controls, full-capacity restore, Autonomous mode, and queue depth.",
     ["Initial state creates 60 RuntimeNode objects. Each starts active with an empty task list and temperature 30. The initial fill attempts four tasks per node, respecting canAccept().", "The effect fetches /api/tasks and then starts a one-second interval. Each tick updates temperature, removes completed or terminated tasks, and optionally performs intake or consolidation.", "makeTask() chooses a duration. Heavy categories (Machine learning, LLM model execution, Gaming) have a minimum duration of 80 seconds; other categories can run for 10 seconds or longer up to the configured upper bound.", "pickTemplate() samples a desired thermal band with a 10% low, 60% medium, and 30% high distribution when matching templates exist. It then returns a template from the catalog.", "metrics() sums CPU, GPU, power, and resource memory, then normalizes RAM by 1200 MB units, VRAM by 600 MB units, and power by 192 W units so the UI can represent percentage-like values.", "canAccept() uses an 80% headroom threshold for CPU, RAM, GPU, VRAM, and power and applies a thermal admission check. This is a safety guard, not a learned prediction.", "temperatureTarget() cools offline nodes toward zero and active empty nodes toward 30. Busy nodes target a bounded value derived from average task temperature and task count.", "Normal mode adds at most one task per node per intake cycle when the node has fewer than 18 tasks, then attempts capacity consolidation.", "Autonomous mode preserves current tasks in a pending queue, stops active nodes sequentially, waits until the room is ready, then dispatches queued work to fitting nodes.", "If no active node can accept queued work, dispatchAutonomousWork() starts the next offline node. The recent autonomous consolidation fix then moves compatible work together and schedules empty nodes for shutdown so green lights do not accumulate without work.", "The intake cadence is approximately 10 seconds. A one-second tick checks elapsed wall time rather than creating a task every second." ]),
    ("8. API and sponsor technology", "The project uses APIs, but they are internal APIs rather than an external vendor integration. This distinction matters in a technical presentation: the prototype demonstrates clean boundaries without claiming a sponsor integration that does not exist.",
     ["GET /api/tasks reads agenticwork/tasks_preview.csv with node:fs/promises, parses the rows, maps them to TaskTemplate records, and returns NextResponse.json().", "GET /api/powerwall reads attached_assets/tesla_powerwall_2_1786387109575.glb and returns the binary bytes with Content-Type model/gltf-binary. GLTFLoader consumes the route from the browser.", "There is no database or authentication API. The login page sets a localStorage flag equin... (the exact key is equinox-auth in the current code) and the dashboard redirects to /login if the flag is absent. This is a demo gate, not production identity.", "No sponsor technology or third-party telemetry API is currently attached. The current technology story is open-source infrastructure: Next.js, React, TypeScript, Three.js, Framer Motion, and browser APIs.", "A real integration can replace /api/tasks with a queue adapter and /api/powerwall with a static asset/CDN path without forcing the dashboard to understand the upstream source." ]),
    ("9. Three.js visualization", "The overview scene is created only while the Overview page is visible. Three.js builds the camera, lights, floor, grid, walls, power lines, rack clones, labels, Powerwall models, and animated power pulses.",
     ["GLTFLoader loads /data_center_rack.glb. The source model is cloned 60 times and arranged in six columns and ten rows.", "Each rack gets a CanvasTexture label showing the node ID and state color. Offline nodes receive a red cross, while active states use green, yellow, or red.", "The scene stores rack roots and uses a Raycaster to map pointer movement or clicks back to a node ID. That ID selects the matching NodeTelemetry panel.", "Power lines are built from CatmullRomCurve3 paths and TubeGeometry. Animated particles move along the curves to imply energy flow.", "The Powerwall route is loaded as a second GLB and mounted next to the supply path. This is a visual demonstration of energy infrastructure, not an electrical simulation.", "The effect cleans up the animation frame, ResizeObserver, controls, renderer, geometries, and materials when the page changes. WebGL context warnings can occur in browser sandboxes without a GPU and do not imply that the route or simulation is broken." ]),
    ("10. UI and UX architecture", "The user journey is deliberately layered. The landing page explains the proposition, About gives context and people, Login gates the demo, and Dashboard moves from room-level visibility to cluster-level operations and node-level detail.",
     ["Overview: the 3D room is the primary visual anchor. Toolbar actions expose Autonomous mode and full-capacity restore. Hovering a rack opens compact telemetry.", "Live Workflow: cluster tabs, a cluster task manager, task rows, node telemetry panels, capacity bars, task durations, and transfer/termination actions make the operational data readable.", "Logs: every important simulation event is represented as a time, kind, and message row. This is essential to explain autonomous transitions during a judging demo.", "Visual language: dark navy surfaces, glass panels, purple brand accents, green nominal state, yellow elevated state, red critical state, and muted metadata.", "Typography: Gotham Black is used for heavy display text, Gotham Light for task names and light display accents, and the monospace face for operational labels.", "Accessibility and readability improvements enlarged the small dashboard panel text, side-panel copy, task metadata, and controls without changing the main heading scale." ]),
    ("11. Team collaboration", "Tanishq and Aaroh divide the work by system boundary while sharing the contract between data, APIs, simulation, and presentation.",
     ["Tanishq Giri — backend and frontend credit: application shell, Next.js routes and page composition, dashboard architecture, Three.js scene, telemetry UI, navigation, styles, simulation integration, and product assembly.", "Aaroh Dharmadhikari — task-library and API credit: creation and curation of the 400-entry workload library, API development for delivering task data and visual assets, workload/resource profiling, AI/data-science framing, and validation/documentation support.", "Shared contract: TaskTemplate fields, resource units, eight categories, node capacity bounds, event types, and the meaning of Autonomous mode.", "Aaroh’s API contribution is concrete in the current codebase: /api/tasks is the data boundary that turns the catalog into typed runtime templates, and /api/powerwall is the asset boundary used by the 3D scene.", "The team workflow is strongest when the library can evolve independently of the UI and when the UI can render a deterministic fixture while APIs are being developed." ]),
    ("12. Open-source availability and contribution", "The current repository does not include an explicit LICENSE file. That means the code should not be described as already open source simply because it uses open-source dependencies.",
     ["Before public release, add a root LICENSE file and a THIRD_PARTY_NOTICES or dependency-policy section. MIT is a sensible default for the application code, subject to the team’s decision.", "Document the attached GLB models, profile images, fonts, and any generated CSV provenance separately. Application-code licensing does not automatically grant rights to every asset.", "An open-source contribution path should publish the task schema, environment setup, demo instructions, simulation assumptions, and examples for adding a new resource category.", "Good first contributions include: configurable topology, pluggable task sources, placement scoring experiments, additional telemetry adapters, accessibility improvements, test fixtures, and scheduler integrations.", "Label the current system as a prototype. Contributors should know which parts are a local browser simulation and which interfaces are intended to become production adapters." ]),
    ("13. Testing and verification", "The project is a live visualization, so verification has two layers: static/build correctness and behavioral inspection in the browser.",
     ["Build check: npm run build should compile the Next.js application and catch TypeScript/import issues.", "Route checks: GET /api/tasks should return an array of task templates with numeric CPU, RAM, GPU, VRAM, power, and temperature fields; GET /api/powerwall should return a GLB binary response.", "UI check: / and /about should load, /login should set the demo gate, and /dashboard should redirect unauthenticated users then render after the local gate is active.", "Behavior check: toggle Autonomous, observe sequential shutdown, watch queued work activate capacity, wait for tasks to complete, and confirm empty nodes consolidate rather than remain green indefinitely.", "Visual check: inspect overview rack labels, hover telemetry, Live Workflow readability, cluster controls, task rows, and Logs. A no-GPU browser may log WebGL warnings while the fallback path remains valid.", "The current app is a demo-first codebase and does not yet contain a formal unit/e2e test suite. The most valuable next tests are pure tests for canAccept, makeTask bounds, task parsing, and autonomous dispatch." ]),
    ("14. Limitations and roadmap", "The honest next step is to graduate the explanatory prototype into a configurable, testable orchestration model without losing the clarity of the demo.",
     ["Resolve topology configuration: the history brief says 30 nodes, while the current dashboard renders 60 racks. Make node count, cluster count, and nodes-per-cluster environment/config values.", "Move mutable simulation state into a dedicated engine or worker if the simulation grows. The current browser hook is appropriate for a demo but not for a long-running control plane.", "Add a real event stream or queue adapter. The current 10-second intake is a deliberate local simulation cadence, not a production workload arrival guarantee.", "Add a scoring interface so a future ML model can rank candidate nodes using historical telemetry, while canAccept() remains a hard safety boundary.", "Add persistent logs, deterministic seeds for repeatable demos, unit tests, and performance instrumentation.", "Replace localStorage demo authentication with a real identity layer before exposing the dashboard to real users.", "Add an explicit license, asset provenance, contributor guide, and deployment/security documentation before public release." ]),
    ("15. Final defense summary", "Equinox is memorable because it makes an invisible systems problem visible without pretending the prototype is already a production data-center controller.",
     ["Problem: mixed workloads are placed without enough awareness of energy, thermal pressure, and idle capacity.", "Solution: a continuous, resource-aware browser simulation that accepts tasks, tracks finite lifetimes, migrates or consolidates work, powers capacity down, and wakes it when demand returns.", "Technical proof: a real Next.js + React application, typed internal APIs, CSV-backed task library, Three.js rack room, live telemetry, autonomous state machine, and event log.", "Team proof: Tanishq owns the full-stack and dashboard experience; Aaroh owns task-library creation, API development, workload/data-science framing, and supporting validation.", "Open-source answer: dependencies are open source, but the repository needs an explicit application license and asset review before it can honestly be called an open-source release.", "Wow answer: flip Autonomous mode and watch the room go dark for a reason, then come alive for a reason." ]),
]


def build_docs():
    md = ["# Equinox Technical Documentation", "", "> Long-form engineering reference for the autonomous energy-aware workload orchestration prototype.", "", "## How to use this document", "", "Every section begins with a **Fast summary** so the team can prepare quickly, followed by the deeper explanation and the code-level details.", ""]
    for title, summary, points in SECTIONS:
        md += [f"## {title}", "", f"**Fast summary:** {summary}", ""]
        for p in points:
            md += [f"- {p}"]
        md += [""]
    md += ["## Appendix: file-level reference", "", "| File | Responsibility |", "|---|---|", "| `app/page.tsx` | Equinox landing page, hero, product framing, feature cards, navigation |", "| `app/about/page.tsx` | Product brief, reference bounds, team profiles |", "| `app/dashboard/page.tsx` | Dashboard shell, Three.js visualization, telemetry panels, cluster operations |", "| `app/dashboard/simulation.ts` | Task lifecycle, telemetry math, placement, migration, autonomous control, power transitions |", "| `app/api/tasks/route.ts` | CSV task-library parser and JSON API |", "| `app/api/powerwall/route.ts` | GLB asset API response |", "| `app/globals.css` | Global page styles and base visual language |", "| `app/overrides.css` | Dashboard/about/font/readability overrides |", "| `agenticwork/tasks_preview.csv` | 400-entry workload catalog |", "| `public/data_center_rack.glb` | Rack model used by Three.js |", "| `attached_assets/tesla_powerwall_2_1786387109575.glb` | Powerwall model served through `/api/powerwall` |", "| `package.json` | Next, React, Three.js, Framer Motion, TypeScript dependencies and run scripts |", "| `tsconfig.json` | TypeScript configuration and reference-build exclusion |", "", "## Appendix: demo script", "", "1. Open the landing page and explain the physical-compute framing.", "2. Use About to introduce the two-person team and the 400-task workload library.", "3. Enter Login, then Dashboard → Overview.", "4. Hover a rack to show CPU, GPU, RAM, VRAM, power, temperature, and active processes.", "5. Open Live Workflow and select a cluster to show task-level operations.", "6. Return to Overview and enable Autonomous mode.", "7. Narrate the sequence: preserve work, drain, power down, slow intake, activate capacity, dispatch, complete, consolidate.", "8. Open Logs to prove that each state transition is observable.", "9. Close with the open-source answer: dependencies are open, project licensing and asset review are next.", ""]
    md_path = OUT / "Equinox_Technical_Documentation.md"
    md_path.write_text("\n".join(md), encoding="utf-8")

    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = DInches(0.7); sec.bottom_margin = DInches(0.7)
    sec.left_margin = DInches(0.8); sec.right_margin = DInches(0.8)
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"; styles["Normal"].font.size = DPt(10); styles["Normal"].font.color.rgb = DRGBColor(45, 51, 65)
    for style_name, size, color in [("Title", 30, (74, 35, 123)), ("Heading 1", 20, (74, 35, 123)), ("Heading 2", 15, (30, 105, 100))]:
        st = styles[style_name]; st.font.name = "Aptos Display"; st.font.size = DPt(size); st.font.bold = True; st.font.color.rgb = DRGBColor(*color)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("EQUINOX"); r.bold = True; r.font.size = DPt(34); r.font.color.rgb = DRGBColor(110, 45, 170)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Technical Documentation"); r.bold = True; r.font.size = DPt(24)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Autonomous energy-aware workload orchestration prototype\n").italic = True
    p.add_run("Prepared for technical presentation, judging, and implementation handoff").font.size = DPt(11)
    doc.add_page_break()
    doc.add_heading("How to use this document", level=1)
    doc.add_paragraph("Every section begins with a Fast summary. Read those summaries first if the team is preparing under time pressure; use the detailed paragraphs and bullets when a judge, teammate, or future contributor asks how the system actually works.")
    doc.add_heading("Table of contents", level=1)
    for title, _, _ in SECTIONS:
        doc.add_paragraph(title, style="List Bullet")
    for title, summary, points in SECTIONS:
        doc.add_page_break()
        doc.add_heading(title, level=1)
        p = doc.add_paragraph()
        r = p.add_run("FAST SUMMARY  "); r.bold = True; r.font.color.rgb = DRGBColor(110, 45, 170)
        p.add_run(summary)
        doc.add_paragraph("Detailed notes", style="Heading 2")
        for item in points:
            doc.add_paragraph(item, style="List Bullet")
        if title == "7. Runtime simulation and control logic":
            doc.add_paragraph("Control-loop walkthrough", style="Heading 2")
            table = doc.add_table(rows=1, cols=3); table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = "Light Shading Accent 1"
            for cell, text in zip(table.rows[0].cells, ["Phase", "Code concept", "What the operator sees"]):
                cell.text = text; cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            rows = [
                ("Tick", "one-second interval", "telemetry and temperatures change"),
                ("Intake", "elapsed time >= 10s", "new task or queue event"),
                ("Admission", "canAccept()", "task lands only with headroom"),
                ("Completion", "remaining <= 0", "task disappears and log is written"),
                ("Autonomous", "dispatch + start", "offline rack wakes for demand"),
                ("Consolidation", "moveTasks + schedulePowerDown", "empty rack powers down"),
            ]
            for row in rows:
                cells = table.add_row().cells
                for cell, text in zip(cells, row): cell.text = text
        if title == "5. Codebase map":
            doc.add_paragraph("File ownership map", style="Heading 2")
            table = doc.add_table(rows=1, cols=2); table.style = "Light Shading Accent 1"
            table.rows[0].cells[0].text = "File"; table.rows[0].cells[1].text = "Role"
            file_rows = [
                ("app/page.tsx", "Landing page and product framing"),
                ("app/about/page.tsx", "System brief and people"),
                ("app/dashboard/page.tsx", "Dashboard and Three.js visualization"),
                ("app/dashboard/simulation.ts", "Runtime state and control loop"),
                ("app/api/tasks/route.ts", "Task catalog API"),
                ("app/api/powerwall/route.ts", "GLB asset API"),
                ("app/globals.css + app/overrides.css", "Visual system and dashboard styling"),
                ("agenticwork/tasks_preview.csv", "400 workload templates"),
            ]
            for file_name, role in file_rows:
                cells = table.add_row().cells; cells[0].text = file_name; cells[1].text = role
    doc.add_page_break()
    doc.add_heading("Appendix: demo script", level=1)
    for step in ["Open the landing page and explain the physical-compute framing.", "Use About to introduce the team and library.", "Enter Login, then Dashboard → Overview.", "Hover a rack and read the telemetry.", "Open Live Workflow for cluster/task detail.", "Enable Autonomous mode and narrate drain → lights out → demand → wake → dispatch → consolidate.", "Open Logs to prove the transitions.", "Answer the open-source question honestly: dependencies are open, project licensing and asset review are next."]:
        doc.add_paragraph(step, style="List Number")
    doc_path = OUT / "Equinox_Technical_Documentation.docx"
    doc.save(str(doc_path))
    return md_path, doc_path


if __name__ == "__main__":
    deck = build_deck()
    md, doc = build_docs()
    print(deck)
    print(md)
    print(doc)